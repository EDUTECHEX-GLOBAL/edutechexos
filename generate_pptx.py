"""
EduTechExOS PowerPoint Presentation Generator
Run: pip install python-pptx && python generate_pptx.py
Outputs: EduTechExOS_Platform_Overview.pptx + EduTechExOS_Speech.txt
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x19, 0x1E, 0x2F)
INDIGO = RGBColor(0x3E, 0x4A, 0x89)
PURPLE = RGBColor(0x4F, 0x46, 0xE5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xA5, 0xB4, 0xFC)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=Pt(18), bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def accent_bar(slide, color=PURPLE):
    box(slide, 0.5, 0.35, 0.06, 0.55, fill_color=color)

def slide_number(slide, n):
    txt(slide, str(n), 12.6, 7.1, 0.6, 0.35, size=Pt(10), color=LIGHT, align=PP_ALIGN.RIGHT)

def divider(slide, top, color=INDIGO):
    box(slide, 0.5, top, 12.33, 0.03, fill_color=color)

def heading(slide, title, subtitle=None, n=None):
    accent_bar(slide)
    txt(slide, title, 0.75, 0.25, 11.5, 0.6, size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.75, 0.85, 11.5, 0.4, size=Pt(14), color=LIGHT, italic=True)
    divider(slide, 1.25)
    if n:
        slide_number(slide, n)

def bullet_column(slide, items, left, top, width=5.5, size=Pt(13), color=WHITE, spacing=0.38):
    for i, item in enumerate(items):
        txt(slide, f"▸  {item}", left, top + i * spacing, width, 0.35, size=size, color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
box(s, 0, 2.5, 13.33, 2.8, fill_color=RGBColor(0x31, 0x2E, 0x81))
txt(s, "EduTechExOS", 0.5, 1.0, 12.33, 1.1, size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "The All-in-One Operating System for EduTechEx Teams", 0.5, 2.1, 12.33, 0.6,
    size=Pt(20), color=LIGHT, align=PP_ALIGN.CENTER)
txt(s, "Real-time Messaging  ·  AI Copilot  ·  Meetings  ·  Kanban  ·  Attendance  ·  Admin",
    0.5, 2.85, 12.33, 0.55, size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
txt(s, f"Presented by EduTechEx  |  {datetime.date.today().strftime('%B %Y')}",
    0.5, 6.8, 12.33, 0.4, size=Pt(11), color=LIGHT, align=PP_ALIGN.CENTER)
slide_number(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem & Why EduTechExOS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "The Problem EduTechExOS Solves", "Why we needed to build our own platform", n=2)

problems = [
    "Slack/Teams = expensive SaaS, no org-specific control",
    "Google Meet is generic, no channel context",
    "Attendance tracked in spreadsheets (error-prone)",
    "Leave approvals via email chains (slow, untracked)",
    "AI tools disconnected from team conversations",
    "No single source of truth for tasks, wiki & comms",
]
solutions = [
    "Self-hosted, zero per-seat SaaS cost",
    "Jitsi meetings embedded in channel context",
    "Auto attendance from login/logout events",
    "Leave calendar + admin approval in one click",
    "Gemini AI reads channel history to answer questions",
    "Kanban + Wiki + Messaging in one unified dashboard",
]
txt(s, "BEFORE (Pain Points)", 0.5, 1.4, 5.8, 0.4, size=Pt(13), bold=True, color=RED)
txt(s, "AFTER (EduTechExOS)", 7.0, 1.4, 5.8, 0.4, size=Pt(13), bold=True, color=GREEN)
for i, (p, sol) in enumerate(zip(problems, solutions)):
    txt(s, f"✗  {p}", 0.5, 1.85 + i*0.73, 5.8, 0.65, size=Pt(11.5), color=RGBColor(0xFC,0xA5,0xA5))
    txt(s, f"✓  {sol}", 7.0, 1.85 + i*0.73, 5.8, 0.65, size=Pt(11.5), color=RGBColor(0x6E,0xE7,0xB7))
box(s, 6.65, 1.3, 0.04, 5.8, fill_color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture & Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Platform Architecture & Tech Stack", "How EduTechExOS is built under the hood", n=3)

layers = [
    ("FRONTEND", "Next.js 15  ·  React 19  ·  TypeScript  ·  Tailwind CSS  ·  Framer Motion", PURPLE),
    ("STATE",    "Zustand (persist)  ·  Socket.IO-client  ·  AES-256 message encryption",     INDIGO),
    ("BACKEND",  "Node.js  ·  Express.js 4  ·  Socket.IO 4  ·  JWT Auth  ·  bcrypt",          GREEN),
    ("DATABASE", "MongoDB + Mongoose  ·  20+ models  ·  GridFS-ready",                        AMBER),
    ("AI",       "Google Gemini (primary)  ·  OpenAI GPT-4o-mini (fallback)  ·  AI SDK",      RGBColor(0xEC,0x48,0x99)),
    ("STORAGE",  "Cloudinary (media)  ·  Backblaze B2  ·  Cloudflare R2  ·  Base64/MongoDB",  RGBColor(0x06,0xB6,0xD4)),
    ("DEPLOY",   "Frontend → Vercel  ·  Backend → Render  ·  Real-time via WebSockets",        RGBColor(0x84,0xCC,0x16)),
]
for i, (label, detail, col) in enumerate(layers):
    top = 1.45 + i * 0.72
    box(s, 0.5, top, 1.5, 0.55, fill_color=col)
    txt(s, label, 0.5, top+0.12, 1.5, 0.35, size=Pt(10), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, detail, 2.2, top+0.1, 10.4, 0.4, size=Pt(12), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Landing Page
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Landing Page & Public Presence", "Route: /  (public-facing marketing + feature overview)", n=4)

features = [
    ("Project Channels",      "Org-wide chat rooms per team/subject, always in sync"),
    ("Embedded AI Agent",     "Gemini-powered assistant trained on your workspace data"),
    ("Auto Task Extraction",  "AI reads messages and creates Kanban cards automatically"),
    ("Morning Digest",        "Daily AI-generated briefing from overnight activity"),
    ("Attendance Tracking",   "Login/logout events auto-logged; calendar view for HR"),
    ("Broadcast & Alerts",    "Org-wide announcements from admin to all members"),
    ("Org Knowledge Base",    "Collaborative wiki with rich-text editor (Tiptap)"),
    ("Activity Monitoring",   "Per-user session logs visible to admins in real time"),
    ("Member Onboarding",     "Email-invite system → guided signup → instant access"),
]
for i, (feat, desc) in enumerate(features):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    left = 0.5 if col == 0 else 6.9
    top  = 1.45 + row * 1.0
    box(s, left, top, 5.9, 0.85, fill_color=RGBColor(0x2D,0x31,0x56), line_color=INDIGO, line_width=Pt(1))
    txt(s, feat, left+0.15, top+0.05, 5.6, 0.3, size=Pt(12), bold=True, color=LIGHT)
    txt(s, desc, left+0.15, top+0.37, 5.6, 0.42, size=Pt(10.5), color=RGBColor(0xC7,0xD2,0xFE))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Auth, Invite System & Security
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Authentication, Invite System & Security", "Routes: /sign-up-login-screen  ·  /invite", n=5)

flows = [
    ("Admin Invite Flow",
     "Admin generates a unique token → email sent → recipient opens /invite link → fills name + password → account created"),
    ("Login",
     "Email + password → bcrypt compare → JWT issued (stored in localStorage) → dashboard loads → Socket.IO connects"),
    ("Access Request",
     "Users without an invite can submit an Access Request → admin sees pending list → approve/decline"),
    ("Security Layers",
     "AES-256 encryption on every message body in MongoDB · JWT expiry · bcrypt hashing · invite tokens are single-use"),
    ("Change Password",
     "Settings → Security tab → old password verified → bcrypt re-hash → token refreshed"),
    ("Removed Members",
     "Admin removes user → account deactivated · stored in RemovedMember collection · re-apply flow available"),
]
for i, (title, body) in enumerate(flows):
    col = i % 2; row = i // 2
    left = 0.5 if col == 0 else 6.9
    top  = 1.45 + row * 1.85
    box(s, left, top, 5.9, 1.7, fill_color=RGBColor(0x1E,0x22,0x40), line_color=PURPLE, line_width=Pt(1))
    txt(s, title, left+0.15, top+0.1, 5.6, 0.4, size=Pt(13), bold=True, color=PURPLE)
    txt(s, body,  left+0.15, top+0.5, 5.6, 1.1, size=Pt(10.5), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Dashboard Overview
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "The Main Dashboard — Overview", "Route: /dashboard  ·  Central hub for all team activity", n=6)

zones = [
    ("Left Sidebar",    "Channel list, DM list, member presence indicators, workspace switcher"),
    ("Top Bar",         "Active channel name, search, meeting button, notification bell, user avatar"),
    ("Message Feed",    "Real-time chat with reactions, replies, pins, bookmarks, file previews"),
    ("Right Rail",      "AI Copilot, Kanban board, Wiki, Leave panel, Insights, Attendance, Search"),
    ("Settings Panel",  "6 tabs: Profile · Appearance · Notifications · Meeting · Security · Privacy"),
    ("Voice/Video",     "In-line voice notes, screen recording, instant Jitsi meetings"),
]
txt(s, "DASHBOARD ZONES", 0.5, 1.4, 12.33, 0.4, size=Pt(12), bold=True, color=LIGHT)
for i, (zone, desc) in enumerate(zones):
    top = 1.85 + i * 0.83
    box(s, 0.5, top, 2.2, 0.7, fill_color=PURPLE)
    txt(s, zone, 0.55, top+0.18, 2.1, 0.4, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, 2.85, top+0.15, 9.9, 0.5, size=Pt(12), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Real-time Messaging
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Real-time Messaging — Feature Deep Dive", "Socket.IO · AES-256 encrypted · instant delivery", n=7)

left_items = [
    "Send rich text messages (Tiptap editor)",
    "Emoji reactions (click any message)",
    "Thread replies per message",
    "Pin a message (private to you only)",
    "Bookmark messages for personal reference",
    "@mention teammates → triggers notification",
    "Real-time typing indicators",
    "Presence: Online / Away / Do Not Disturb",
]
right_items = [
    "File attachments: images, video, audio, PDF/DOC/ZIP",
    "Images/video stored on Cloudinary (CDN)",
    "Raw files (PDF etc.) stored as base64 in MongoDB",
    "Screen recording sent directly to the channel",
    "Voice notes recorded in-browser, shared inline",
    "Message search with full-text keyword filter",
    "Deduplication via clientId → server echo",
    "Optimistic UI: message shows instantly, confirmed by socket",
]
txt(s, "Interaction Features", 0.5, 1.4, 5.8, 0.4, size=Pt(13), bold=True, color=PURPLE)
txt(s, "File & Media", 7.0, 1.4, 5.8, 0.4, size=Pt(13), bold=True, color=GREEN)
bullet_column(s, left_items, 0.5, 1.85, width=5.9, size=Pt(11.5))
bullet_column(s, right_items, 7.0, 1.85, width=5.9, size=Pt(11.5))
box(s, 6.65, 1.3, 0.04, 5.8, fill_color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Voice, Recording & Files
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Voice Notes, Screen Recording & File Sharing", "In-browser media capture — no external app needed", n=8)

cards = [
    ("Voice Notes",
     "Click microphone icon in MessageInput → MediaRecorder captures audio in-browser → "
     "uploaded to Cloudinary (audio/*) → audio player embedded in message feed → "
     "all channel members can play it back instantly."),
    ("Screen Recording",
     "Click record button → browser prompts screen + audio permission → "
     "MediaRecorder records the stream → on stop: file auto-attached to message input → "
     "user sends it → stored on Cloudinary → sender AND all members see the video inline."),
    ("File Attachments",
     "Drag & drop or click paperclip → smartUpload() routes by type: "
     "images/video/audio → Cloudinary · PDF/DOC/ZIP → base64 MongoDB (max 5 MB) · "
     "media files max 10 MB. Download links shown in message feed."),
    ("Permissions Handling",
     "Pre-flight microphone permission check before any recording attempt. "
     "If browser has blocked mic: actionable toast with exact steps "
     "(click lock icon → Microphone → Allow → refresh). "
     "Screen-share cancel vs audio-block produce distinct error messages."),
]
for i, (title, body) in enumerate(cards):
    col = i % 2; row = i // 2
    left = 0.5 if col == 0 else 6.9
    top  = 1.45 + row * 2.6
    box(s, left, top, 5.9, 2.4, fill_color=RGBColor(0x1E,0x22,0x40), line_color=INDIGO, line_width=Pt(1))
    txt(s, title, left+0.15, top+0.1, 5.6, 0.45, size=Pt(14), bold=True, color=LIGHT)
    txt(s, body,  left+0.15, top+0.55, 5.6, 1.75, size=Pt(10.5), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Meeting Management
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Meeting Management", "Instant Jitsi meetings · Scheduled sessions · /meeting/[code] room", n=9)

txt(s, "INSTANT MEETINGS", 0.5, 1.4, 5.8, 0.4, size=Pt(13), bold=True, color=GREEN)
instant = [
    "Click 'Start Meeting' in any channel",
    "System generates: edutechexos-{channel}-{timestamp}",
    "Jitsi URL: https://meet.jit.si/{roomSlug}",
    "Link stored in meeting notification with joinLink",
    "All channel members receive notification",
    "Any member clicks 'Join Meeting' button in alerts",
    "Everyone joins the SAME Jitsi room",
    "Free, open-source, no account required",
]
txt(s, "SCHEDULED MEETINGS", 7.0, 1.4, 5.8, 0.4, size=Pt(13), bold=True, color=AMBER)
scheduled = [
    "Admin or user schedules via Meeting Settings tab",
    "Date, time, channel, description set",
    "Stored in MeetingRequest collection",
    "Notification broadcast to all invited members",
    "joinLink included in notification for one-click join",
    "/meeting/[code] route: dedicated room page",
    "MeetingAccess model controls who can enter",
    "Participant list, join/leave events tracked",
]
bullet_column(s, instant, 0.5, 1.85)
bullet_column(s, scheduled, 7.0, 1.85)
box(s, 6.65, 1.3, 0.04, 5.8, fill_color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — AI Copilot & Automation
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "AI Copilot & Smart Automation", "Google Gemini primary · OpenAI GPT-4o-mini fallback", n=10)

ai_features = [
    ("Conversational Q&A",
     "Ask the AI Copilot anything about your workspace, team, channels, or tasks. "
     "Gemini reads channel history and member context before answering."),
    ("Auto Task Extraction",
     "AI monitors messages for action items and automatically creates Kanban cards "
     "in the correct column — no manual ticket creation needed."),
    ("Morning Digest",
     "Scheduled job generates a daily AI briefing summarising overnight messages, "
     "new tasks, and attendance patterns — sent to every member at start of day."),
    ("Smart Search",
     "Natural-language search across messages, wiki pages, and Kanban cards. "
     "AI ranks results by relevance, not just keyword match."),
    ("Availability Insights",
     "Admin Availability tab uses AI to surface patterns: who is consistently "
     "away, peak active hours, suggested meeting windows."),
    ("Fallback Logic",
     "If Gemini API fails or rate-limits: automatic fallback to OpenAI GPT-4o-mini. "
     "User sees no interruption — fully transparent failover."),
]
for i, (title, body) in enumerate(ai_features):
    col = i % 2; row = i // 2
    left = 0.5 if col == 0 else 6.9
    top  = 1.45 + row * 1.85
    box(s, left, top, 5.9, 1.7, fill_color=RGBColor(0x1E,0x18,0x3B), line_color=PURPLE, line_width=Pt(1))
    txt(s, title, left+0.15, top+0.1, 5.6, 0.4, size=Pt(13), bold=True, color=LIGHT)
    txt(s, body,  left+0.15, top+0.5, 5.6, 1.1, size=Pt(10.5), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Kanban Board & Wiki
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Workspace Tools — Kanban Board & Wiki", "Right-rail panels · KanbanTask + WikiPage models", n=11)

txt(s, "KANBAN BOARD", 0.5, 1.4, 5.8, 0.4, size=Pt(14), bold=True, color=AMBER)
kanban = [
    "4 columns: Backlog → In Progress → Review → Done",
    "Create cards manually or via AI extraction",
    "Drag & drop cards between columns",
    "Assign cards to team members",
    "Set due dates and priority levels",
    "Card detail modal: description, comments, attachments",
    "Per-channel or workspace-wide view toggle",
    "Real-time sync: card moves visible to all via Socket.IO",
]
txt(s, "ORG KNOWLEDGE BASE (WIKI)", 7.0, 1.4, 5.8, 0.4, size=Pt(14), bold=True, color=GREEN)
wiki = [
    "Rich text editor powered by Tiptap",
    "Supports headings, bold, lists, code blocks, links",
    "Create, edit, delete wiki pages",
    "Pages stored in WikiPage MongoDB collection",
    "Per-page version: author + timestamp visible",
    "All members can read; admins/creators can edit",
    "Search integration: wiki pages appear in global search",
    "Sidebar navigation lists all pages alphabetically",
]
bullet_column(s, kanban, 0.5, 1.85)
bullet_column(s, wiki, 7.0, 1.85)
box(s, 6.65, 1.3, 0.04, 5.8, fill_color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Leave Management & Attendance
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Leave Management & Attendance Tracking", "Leave · LoginEvent · ActivitySession · AuditLog", n=12)

txt(s, "LEAVE MANAGEMENT", 0.5, 1.4, 5.8, 0.4, size=Pt(14), bold=True, color=PURPLE)
leave = [
    "Member submits leave request from dashboard",
    "Sets: type (casual/sick/earned), dates, reason",
    "Admin sees pending requests in Admin → Leaves tab",
    "One-click Approve / Decline with comment",
    "Approved leaves appear in Leave Calendar (visual)",
    "Members can see their own leave history",
    "Leave balance tracked per member per type",
    "Admin can manually adjust leave balances",
]
txt(s, "ATTENDANCE TRACKING", 7.0, 1.4, 5.8, 0.4, size=Pt(14), bold=True, color=GREEN)
attendance = [
    "Auto-logged on every login event (LoginEvent model)",
    "Logout / session-end also recorded",
    "ActivitySession tracks active usage duration",
    "UserAttendanceCalendar: heatmap view per user",
    "LoginTrackerCalendar: month view with daily status",
    "Admin sees ALL members' attendance in one view",
    "Late arrivals, absences flagged automatically",
    "Export-ready data for HR/payroll integration",
]
bullet_column(s, leave, 0.5, 1.85)
bullet_column(s, attendance, 7.0, 1.85)
box(s, 6.65, 1.3, 0.04, 5.8, fill_color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Admin Panel
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Admin Panel — Complete Control Center", "Route: /admin  ·  12 management tabs", n=13)

tabs = [
    ("People",          "View all members, roles, status; remove or promote users"),
    ("Requests",        "Approve / decline incoming access requests"),
    ("Invites",         "Generate invite tokens; view pending/accepted invites"),
    ("Channels",        "Create, rename, archive workspace channels"),
    ("Broadcast",       "Send org-wide announcements to all members at once"),
    ("Activity",        "Real-time activity log; /admin/activity for drill-down"),
    ("Attendance",      "Full attendance grid: all members x all days"),
    ("Desktop",         "AWActivity model: track which apps members use (opt-in)"),
    ("Availability",    "AdminAvailability: set office hours, view who's online now"),
    ("Leaves",          "Manage all leave requests; approve/decline; view calendar"),
    ("Leave Calendar",  "Visual month-grid of approved leaves across all members"),
    ("Audit",           "AuditLog: every admin action timestamped for compliance"),
]
cols = [tabs[:6], tabs[6:]]
for ci, col_tabs in enumerate(cols):
    left = 0.5 if ci == 0 else 6.9
    for ri, (tab, desc) in enumerate(col_tabs):
        top = 1.4 + ri * 0.97
        box(s, left, top, 5.9, 0.82, fill_color=RGBColor(0x1E,0x22,0x40), line_color=INDIGO, line_width=Pt(1))
        txt(s, tab, left+0.15, top+0.07, 2.0, 0.35, size=Pt(12), bold=True, color=PURPLE)
        txt(s, desc, left+2.2, top+0.07, 3.5, 0.65, size=Pt(10.5), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Notifications, Alerts & Presence
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Notifications, Alerts & Presence System", "Socket.IO events · per-user filtering · real-time bell", n=14)

notif_types = [
    ("Reply",    "#3E4A89", "Someone replies to your message in any channel"),
    ("Reaction", "#F59E0B", "Emoji reaction added to your message"),
    ("Pin",      "#8B5CF6", "A message in your channel is pinned"),
    ("Task",     "#10B981", "A Kanban card is assigned to you"),
    ("Mention",  "#EF4444", "@mention in any channel message"),
    ("Meeting",  "#10B981", "Meeting started / scheduled — includes Join Meeting button"),
]
txt(s, "NOTIFICATION TYPES", 0.5, 1.4, 12.33, 0.4, size=Pt(13), bold=True, color=LIGHT)
for i, (ntype, hex_color, desc) in enumerate(notif_types):
    col = i % 2; row = i // 2
    left = 0.5 if col == 0 else 6.9
    top  = 1.85 + row * 1.5
    r,g,b = int(hex_color[1:3],16), int(hex_color[3:5],16), int(hex_color[5:7],16)
    color = RGBColor(r,g,b)
    box(s, left, top, 0.6, 1.3, fill_color=color)
    txt(s, ntype, left+0.65, top+0.1, 4.9, 0.4, size=Pt(13), bold=True, color=color)
    txt(s, desc, left+0.65, top+0.5, 4.9, 0.7, size=Pt(11.5), color=WHITE)

txt(s, "PRESENCE SYSTEM", 0.5, 6.2, 12.33, 0.4, size=Pt(12), bold=True, color=LIGHT)
txt(s, "Online (green) · Away (amber, after 5m idle) · Do Not Disturb (red, user-set) · Offline (grey) — updated via Socket.IO heartbeat every 30s",
    0.5, 6.55, 12.33, 0.55, size=Pt(11), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Pros, Cons & Honest Assessment
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Pros, Cons & Honest Assessment", "An objective look at EduTechExOS", n=15)

pros = [
    "Zero per-seat SaaS cost — fully self-hosted",
    "Completely customisable for EduTechEx workflows",
    "AI deeply integrated (not bolted-on)",
    "End-to-end AES-256 message encryption",
    "Jitsi meetings: free, no account, no limits",
    "Real-time everything via Socket.IO",
    "Attendance auto-tracked — zero admin effort",
    "Single app replaces 5+ external tools",
    "Leave, wiki, kanban all in one dashboard",
    "Full data ownership & compliance control",
]
cons = [
    "Render free tier → 50s cold start on first visit",
    "No native mobile app (web only currently)",
    "File size limits: 5 MB raw, 10 MB media",
    "Backblaze B2/R2 storage not yet configured",
    "No end-to-end encrypted voice/video (Jitsi handles it)",
    "Single-tenant: one workspace per deployment",
    "No offline mode — requires internet connection",
    "Limited third-party integrations",
    "Admin panel UX can feel dense for new admins",
    "No calendar sync (Google Calendar / Outlook)",
]
txt(s, "PROS", 0.5, 1.4, 5.8, 0.4, size=Pt(14), bold=True, color=GREEN)
txt(s, "CONS", 7.0, 1.4, 5.8, 0.4, size=Pt(14), bold=True, color=RED)
bullet_column(s, pros, 0.5, 1.85, size=Pt(11), color=RGBColor(0x6E,0xE7,0xB7), spacing=0.48)
bullet_column(s, cons, 7.0, 1.85, size=Pt(11), color=RGBColor(0xFC,0xA5,0xA5), spacing=0.48)
box(s, 6.65, 1.3, 0.04, 5.8, fill_color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Why Only This Application?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Why Only EduTechExOS?", "The unique value proposition for EduTechEx teams", n=16)

reasons = [
    ("Built FOR Us, BY Us",
     "Every feature maps to a real EduTechEx workflow. "
     "No vendor decides what we can or can't do. "
     "We ship a feature request in hours, not months."),
    ("Cost = Zero Ongoing SaaS",
     "Slack Pro: $7.25/seat/month = Rs 7,000+/month for 10 people. "
     "Zoom: $150/month. Notion: $96/month. "
     "EduTechExOS: self-hosted, pay only for compute."),
    ("Data Privacy & Ownership",
     "All messages, files, attendance, and leave data live in "
     "our own MongoDB. No third-party reads our data. "
     "AES-256 encryption on message content at rest."),
    ("AI That Knows Our Context",
     "Unlike generic AI tools, our Gemini Copilot reads our "
     "actual channel history, member list, and task board. "
     "Answers are grounded in real workspace data."),
    ("One Login, Everything",
     "Dashboard → messaging, meetings, AI, tasks, wiki, leave, "
     "attendance — all in one tab. No switching apps. "
     "Members stay in flow instead of context-switching."),
    ("Grows With Us",
     "Open codebase: we can add custom modules anytime. "
     "Mobile app, calendar sync, LMS integration — all on the "
     "roadmap and achievable without any vendor dependency."),
]
for i, (title, body) in enumerate(reasons):
    col = i % 2; row = i // 2
    left = 0.5 if col == 0 else 6.9
    top  = 1.45 + row * 1.85
    box(s, left, top, 5.9, 1.7, fill_color=RGBColor(0x1E,0x22,0x40), line_color=GREEN, line_width=Pt(1))
    txt(s, title, left+0.15, top+0.1, 5.6, 0.4, size=Pt(13), bold=True, color=GREEN)
    txt(s, body,  left+0.15, top+0.5, 5.6, 1.1, size=Pt(10.5), color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Future Roadmap & Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "Future Roadmap & Thank You", "What comes next for EduTechExOS", n=17)

roadmap = [
    ("Q3 2026", "Native mobile app (React Native) — iOS + Android"),
    ("Q3 2026", "Google Calendar / Outlook sync for meetings and leaves"),
    ("Q4 2026", "Backblaze B2 storage fully configured for large file uploads"),
    ("Q4 2026", "Multi-workspace support — one platform, many org units"),
    ("Q1 2027", "LMS integration: link course materials to channels"),
    ("Q1 2027", "AI auto-generated meeting minutes from recordings"),
    ("Q2 2027", "Public API + Zapier integration for external workflows"),
    ("Q2 2027", "End-to-end encrypted DMs (Signal Protocol)"),
]
for i, (quarter, item) in enumerate(roadmap):
    top = 1.45 + i * 0.6
    box(s, 0.5, top, 1.6, 0.5, fill_color=PURPLE)
    txt(s, quarter, 0.5, top+0.1, 1.6, 0.35, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, item, 2.3, top+0.1, 10.5, 0.4, size=Pt(12), color=WHITE)

box(s, 0.5, 6.3, 12.33, 0.9, fill_color=INDIGO)
txt(s, "Thank You — Questions Welcome",
    0.5, 6.35, 12.33, 0.5, size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save PPTX ─────────────────────────────────────────────────────────────────
pptx_path = "EduTechExOS_Platform_Overview.pptx"
prs.save(pptx_path)
print(f"[OK] Saved {pptx_path}")

# ══════════════════════════════════════════════════════════════════════════════
# SPEECH DOCUMENT
# ══════════════════════════════════════════════════════════════════════════════
speech = """EduTechExOS — Complete Presentation Speech
Estimated delivery time: 45-55 minutes at a comfortable pace
============================================================

BEFORE YOU START:
- Open the PowerPoint in presenter view so you see these notes on your screen
- Have a glass of water ready
- Pace yourself: each slide should take 2-4 minutes
- Pause for audience questions after every 3-4 slides if time allows

============================================================
SLIDE 1 - TITLE (2 minutes)
============================================================

Good [morning / afternoon / evening] everyone.

My name is [Your Name], and today I'm going to take you through EduTechExOS,
the all-in-one operating system we've built specifically for EduTechEx.

Before I dive in, I want to set the scene. How many of you have ever felt
frustrated switching between Slack, Google Meet, Notion, a spreadsheet for
attendance, email for leave requests, and maybe a separate task board — all
in the same workday? That context-switching is exhausting, and it's expensive,
both in time and in money.

EduTechExOS was built to eliminate exactly that problem. Everything your team
does — messaging, video calls, task management, leave requests, attendance
tracking, knowledge sharing, and AI assistance — all of it lives in a single
application that we built ourselves and that we fully control.

Today I'll walk you through every page, every feature, how it all works under
the hood, why we chose to build it this way, and what's coming next.
Let's get started.

============================================================
SLIDE 2 - THE PROBLEM (3 minutes)
============================================================

Let me start with the problem.

On the left side of this slide, you'll see the world BEFORE EduTechExOS.

First: tools like Slack and Teams are expensive SaaS products. You pay per
seat, per month, and you have no control over the data they store about your
team. For a growing team, that's thousands of rupees every month — money that
could go toward actual product development.

Second: Google Meet is generic. There's no connection between a meeting and
the channel discussion that triggered it. You have to share links manually,
and those links are not tied to any context.

Third: attendance was tracked in spreadsheets. Every morning someone had to
manually update a sheet. Errors crept in. Data was lost. HR had no reliable
source of truth.

Fourth: leave approvals happened over email. A member sends a request, waits
for an admin to see it, then waits again for approval. There was no single
place to see who was on leave on any given day.

Fifth: AI tools like ChatGPT have no idea who your team members are, what
channels you have, or what tasks are in progress. They're general-purpose,
not workspace-aware.

And sixth: tasks, documentation, messaging, and attendance all lived in
different applications. There was no unified view of team activity.

Now look at the right side. With EduTechExOS, EVERY single one of those
problems is solved. The platform is self-hosted, so cost is just compute.
Meetings are tied to channels. Attendance is auto-logged from login events.
Leave is managed in a dedicated panel with one-click approval. Our AI copilot
reads your actual workspace data. And everything lives in one application.

============================================================
SLIDE 3 - ARCHITECTURE & TECH STACK (4 minutes)
============================================================

Now let's look at how EduTechExOS is built.

I've broken the tech stack into seven layers, each color-coded.

The FRONTEND is built on Next.js 15 with React 19 and TypeScript. We use
Tailwind CSS for styling and Framer Motion for animations. The result is a
fast, server-side-rendered application that loads quickly even on slow
connections.

STATE MANAGEMENT uses Zustand with a persistence layer. This means your
preferences, dark mode setting, and channel data survive page refreshes
without hitting the server again. We also use Socket.IO on the client side
for real-time events, and all message content is encrypted with AES-256
before it's sent or stored.

The BACKEND is a Node.js and Express.js server — straightforward, fast, and
well-understood. Socket.IO handles all real-time communication: typing
indicators, presence updates, new messages, notifications.

The DATABASE is MongoDB with Mongoose. We have over 20 collections — messages,
users, notifications, kanban tasks, wiki pages, leave requests, attendance
logs, audit records, and more. MongoDB's flexible schema was the right choice
for a product that has evolved rapidly.

AI is powered by Google Gemini as the primary model. If Gemini hits a rate
limit or fails, we automatically fall back to OpenAI's GPT-4o-mini. The
user never sees this switch — it's completely transparent.

STORAGE uses a priority chain. We first try Backblaze B2, then Cloudflare R2.
For raw files like PDFs and documents, we skip cloud storage entirely and
store them as base64 directly in MongoDB, which keeps things simple and avoids
upload permission issues. Images, video, and audio go to Cloudinary for CDN
delivery.

Finally, DEPLOYMENT: the frontend is on Vercel — globally distributed, fast
CDN, zero configuration. The backend is on Render's free tier. The one
trade-off there is a 50-second cold start if the server hasn't been accessed
recently, but we handle that with a retry mechanism on the client.

============================================================
SLIDE 4 - LANDING PAGE (3 minutes)
============================================================

The landing page at the root URL is our public-facing marketing page, but it
also serves as a feature showcase for anyone considering EduTechExOS for their
own team.

Let me walk through the nine feature cards on this page.

Project Channels: the foundation of the platform. Every team or subject area
gets its own channel. Messages, files, and meetings are always in context.

Embedded AI Agent: our Gemini-powered assistant is embedded directly in the
dashboard. It reads your channel history, your task board, and your team
roster before answering questions. This is not a generic chatbot.

Auto Task Extraction: when someone says "can you finish the report by Friday"
in a channel message, the AI notices that and creates a Kanban card
automatically. No manual ticket creation needed.

Morning Digest: every morning, every member gets a personalized AI-generated
briefing. It summarizes what happened overnight in their channels, lists their
pending tasks, and highlights any attendance flags.

Attendance Tracking: every login event is recorded. Every logout too. The
system builds a complete picture of who was active, when, and for how long,
automatically, with no manual input.

Broadcast & Alerts: admins can send org-wide announcements with one click.
Every member sees the alert in their notification panel.

Org Knowledge Base: a collaborative wiki with rich text editing. Teams can
document processes, decisions, onboarding guides, and anything else that
needs to be preserved.

Activity Monitoring: admins can see real-time activity logs — which channels
are most active, which members have been online, what actions were taken.

And Member Onboarding: the invite flow is designed so that getting a new
person into the platform takes less than two minutes.

============================================================
SLIDE 5 - AUTH, INVITE SYSTEM & SECURITY (4 minutes)
============================================================

Let's talk about how members get in and how security is handled.

The primary flow is the Admin Invite Flow. An admin navigates to the Invites
tab in the admin panel, enters an email address, and clicks generate. The
system creates a unique invite token and sends an email. The recipient opens
the link — which goes to the /invite route — fills in their name and password,
and their account is created. The invite token is single-use and expires.
This keeps access tightly controlled.

The Login flow uses email and password. The password is checked against a
bcrypt hash stored in MongoDB. If it matches, a JWT token is issued. That
token is stored in localStorage under the key edutechex_token. When a Socket.IO
connection is established, that token is used to authenticate the socket as
well, so only logged-in users can receive real-time events.

We also have an Access Request flow. If someone wants to join but doesn't have
an invite, they can submit a request with their name, email, and a reason.
Admins see these in the Requests tab and can approve or decline with a note.

Security is multi-layered. Every message body is encrypted with AES-256 before
it's written to MongoDB. Even if someone got direct database access, they would
see ciphertext, not readable messages. JWT tokens expire, forcing periodic
re-authentication. Passwords are never stored — only bcrypt hashes. Invite
tokens are one-time-use.

When a member is removed by an admin, their account is deactivated immediately.
Their record is moved to the RemovedMember collection. If they want to return,
they go through the re-apply flow, which requires explicit admin approval.

============================================================
SLIDE 6 - DASHBOARD OVERVIEW (3 minutes)
============================================================

Now we arrive at the heart of the application: the main dashboard.

When you log in, you land on the dashboard at /dashboard. The layout has five
distinct zones.

The Left Sidebar shows your channel list, direct message threads, and a
presence indicator for each member. Online members have a green dot. Away
members have amber. Do Not Disturb shows red. Offline is grey. These states
update via Socket.IO every 30 seconds.

The Top Bar shows the active channel name, a search button, a meeting button
to start an instant Jitsi session, a notification bell with an unread badge,
and your user avatar which opens settings.

The Message Feed takes up the central area. This is where real-time
conversation happens. Messages appear instantly — first as an optimistic local
preview, then replaced by the server-confirmed version when the socket echo
arrives. This gives zero-latency feel while maintaining data consistency.

The Right Rail is a panel that slides open with different tools: the AI
Copilot, the Kanban board, the Wiki, Leave requests, Insights, Attendance
calendar, and the Search panel.

And finally, the Settings Panel is accessible from your avatar. It has six
tabs covering your profile, appearance preferences, notification settings,
meeting configuration, security (change password), and privacy controls.

============================================================
SLIDE 7 - REAL-TIME MESSAGING (4 minutes)
============================================================

Let me go deep on messaging, because it's the most-used feature.

On the interaction side: the message composer uses Tiptap — a full rich text
editor. You can add emoji reactions to any message by clicking it. You can
reply in a thread, keeping long conversations organized. You can pin a message
to your personal pin list — and importantly, pinned messages are private to
you; other members don't see your pins. You can bookmark messages for later
reference — again, personal to you. You can @mention teammates, which triggers
an immediate notification to them.

While you type, other members see a typing indicator at the bottom of the feed.
This is powered by a Socket.IO event that fires on keydown and clears after
two seconds of inactivity.

On the file and media side: you can attach any file type. The system
automatically routes it to the right storage. Images, video, and audio go to
Cloudinary for fast CDN delivery. PDFs, Word documents, ZIP files, and other
raw formats go directly to MongoDB as base64 — this avoids the permission
restrictions that cloud storage providers have on unsigned uploads of non-media
files.

The size limits are five megabytes for raw files and ten megabytes for media.
If a file exceeds that, the user gets a clear error toast before the upload
even starts.

Message deduplication is handled through a clientId system. When you send a
message, it appears immediately with a temporary local ID. The server saves it,
assigns a permanent MongoDB ID, and echoes it back via socket. The client
replaces the local version with the server-confirmed version, matching on
either the permanent ID or the clientId. This is what ensures bookmarks and
pins work correctly — they always reference stable server IDs, never temporary
local ones.

============================================================
SLIDE 8 - VOICE, RECORDING & FILES (3 minutes)
============================================================

Let's talk about the media features.

Voice Notes: click the microphone icon in the message input bar. The browser
requests microphone permission. If granted, MediaRecorder starts capturing
audio. When you stop, the recording is attached to your message input
automatically. You send it like any message. The audio player appears inline
in the message feed. Every channel member can play it back.

Screen Recording: click the record button. The browser shows a screen-picker
dialog — you choose which screen or window to share. MediaRecorder captures
the screen output. When you stop recording, the video file is attached to your
message input. Send it, and everyone sees a video player inline.

An important detail here: the sender can always see their own screen recording.
This was a bug we fixed — previously, only other members could see the video.
The fix was making sure the message comes back through the server echo and is
properly rendered in the sender's feed.

For both voice and screen recording, we've added pre-flight permission checks.
If your browser has blocked microphone access — which can happen after declining
the permission prompt — you'll see an actionable error message: "Click the lock
icon in the address bar, find Microphone, set it to Allow, then refresh." This
replaces the unhelpful generic error message that we had before.

============================================================
SLIDE 9 - MEETING MANAGEMENT (3 minutes)
============================================================

Meetings in EduTechExOS are powered by Jitsi Meet, which is a free,
open-source video conferencing platform. No account required, no per-minute
billing, no participant limits.

For an instant meeting: click the meeting button in the top bar while you're
in a channel. The system generates a room name using the format
edutechexos-{channel-name}-{timestamp}. It then builds a Jitsi URL:
https://meet.jit.si/{roomName}. This link is stored in a meeting notification
and broadcast to all channel members. Anyone who clicks "Join Meeting" in their
notification panel opens the exact same Jitsi URL. Everyone ends up in the
same room.

This is deterministic. The room URL is derived from the channel name and
timestamp, so there's no ambiguity. Every member who receives the notification
clicks the same link and joins the same call.

For scheduled meetings: admins or users navigate to the Meeting tab in
settings, fill in the date, time, description, and channel, and submit. The
system creates a MeetingRequest record in MongoDB and sends notifications to
all invited members. The notification includes the joinLink so members can
join directly from the notification panel.

We also have a dedicated meeting room at /meeting/[code]. This is a full-page
meeting room experience with participant management and access control via the
MeetingAccess model.

============================================================
SLIDE 10 - AI COPILOT & AUTOMATION (4 minutes)
============================================================

AI in EduTechExOS is not a chatbot widget bolted to the side. It's deeply
integrated into how the platform works.

The Conversational AI Copilot in the right rail can answer any question about
your workspace. Ask it "what did the engineering channel discuss this week?"
and it reads the actual message history before responding. Ask it "who is
assigned to the homepage redesign task?" and it reads the Kanban board. This
context-awareness is what makes it genuinely useful rather than generic.

Auto Task Extraction: the AI monitors messages and looks for action items.
"Can you update the API documentation by Thursday" — the AI extracts that,
creates a Kanban card with the assignee and due date, and places it in the
Backlog column. No one has to manually create a ticket.

Morning Digest: a scheduled backend job runs every morning. It calls Gemini
with a prompt that includes overnight message activity, pending tasks, and
attendance flags. Gemini generates a personalized summary for each member. This
lands in their notifications when they start their day.

The AI-powered search goes beyond keyword matching. Type a natural-language
query into the search panel and it finds relevant messages, wiki pages, and
Kanban cards even if the exact words don't match.

Availability Insights in the admin panel use AI to surface patterns. If a
member is consistently offline during supposed working hours, or if there's a
predictable pattern of peak activity, the AI surfaces that for the admin.

And the failover: if Gemini hits a rate limit or returns an error, we
automatically retry with OpenAI's GPT-4o-mini. The switch is invisible to
the user. They might get a slightly different response style, but the function
is identical.

============================================================
SLIDE 11 - KANBAN BOARD & WIKI (3 minutes)
============================================================

Two of the most-loved features in the right rail are the Kanban board and
the Wiki.

The Kanban board has four columns: Backlog, In Progress, Review, and Done.
Cards can be created manually — click add card, fill in title, description,
assignee, and due date. Or they can be created automatically by the AI when
it detects action items in messages.

Drag a card from Backlog to In Progress when work starts. Drag it to Review
when it's ready for feedback. Drag it to Done when it's complete. These moves
are broadcast via Socket.IO, so every team member sees the board update in
real time without refreshing.

The card detail modal lets you add comments, attach files, set priority, and
view the full history of changes. You can filter by assignee or by channel.

The Wiki is a collaborative knowledge base powered by Tiptap. Every page
supports rich formatting: headings, lists, code blocks, bold, italic, and links.
Pages are stored in the WikiPage collection in MongoDB. Any member can create a
new page. Admins and page creators can edit or delete pages. Pages are listed
alphabetically in the sidebar.

The wiki is indexed by the search system, so typing a keyword in the search
panel returns matching wiki pages alongside messages and Kanban cards.

============================================================
SLIDE 12 - LEAVE MANAGEMENT & ATTENDANCE (4 minutes)
============================================================

Leave management and attendance are features that really differentiate
EduTechExOS from a generic messaging tool.

For leave: a member opens the Leave panel in the right rail. They fill in the
leave type — casual, sick, or earned — the start date, end date, and reason.
They submit. The request immediately appears in the Admin panel under the
Leaves tab as "Pending."

The admin sees all pending requests. They can approve with one click, or
decline with a comment explaining why. When a leave is approved, it appears
on the Leave Calendar — a visual month-grid that shows all approved leaves
across all members. Any member can see this calendar and plan their own
requests accordingly.

Admins can also manually adjust leave balances when needed — for example,
carrying over unused casual leave from a previous period.

For attendance: every time a member logs in, a LoginEvent record is created
in MongoDB with the timestamp, IP address, and user agent. When they log out
or their session ends, another event is recorded. ActivitySession models track
how long a member was actively using the application versus just having it
open in the background.

The UserAttendanceCalendar component shows a heatmap view — darker color
means more active time that day. The LoginTrackerCalendar shows a month
view with a green/red/amber indicator for each day: present, absent, or late.

Admins see this for ALL members in the Attendance tab, giving HR a complete
real-time picture of team activity without any manual data entry.

============================================================
SLIDE 13 - ADMIN PANEL (5 minutes)
============================================================

The admin panel at /admin is the command center. Let me walk through all
twelve tabs because each one serves a distinct purpose.

People: view every member of the workspace. See their role, their current
presence status, when they last logged in. Remove a member, which immediately
deactivates their account. Promote someone to admin to grant them panel access.

Requests: the queue of people who submitted access requests without an invite.
Review each request, see their stated reason, and approve or decline. If
approved, they get an invite link automatically.

Invites: generate new invite tokens for specific email addresses. See which
invites have been accepted and which are still pending. Revoke a pending invite
if needed.

Channels: create new channels for new teams or subjects. Rename existing
channels. Archive channels that are no longer active — archived channels are
hidden from the sidebar but messages are preserved.

Broadcast: type an announcement message and send it to every single member
simultaneously. This appears in everyone's notification panel as a broadcast
alert. No email, no separate system — one click reaches the whole org.

Activity: a real-time activity log showing which members are active, what
channels they've been posting in, and how many messages were sent in the last
24 hours. There's also a dedicated /admin/activity route with drill-down
filters.

Attendance: the full attendance grid — every member on one axis, every day on
the other. Color-coded cells give HR an immediate visual overview.

Desktop: the AWActivity tab tracks which desktop applications members are using
during work hours. This is opt-in and uses an activity tracker agent that
members install voluntarily. The data feeds into productivity analysis.

Availability: the AdminAvailability model lets admins set official office hours
and see in real time which members are currently online versus away. AI
surfaces patterns about availability trends.

Leaves: the full leave management interface — all pending, approved, and
declined requests in one view. Filter by date range, member, or leave type.

Leave Calendar: a visual calendar showing all approved leaves as blocks across
the month. Admins use this for resource planning — you can see at a glance
that three people are on leave next Monday before scheduling a critical meeting.

And finally, Audit: every admin action is written to the AuditLog collection
with a timestamp, the admin's identity, and the action taken. Who removed which
member and when. Who approved which leave request. Who created which channel.
This is essential for compliance and accountability.

============================================================
SLIDE 14 - NOTIFICATIONS, ALERTS & PRESENCE (3 minutes)
============================================================

The notification system in EduTechExOS is designed around one principle: only
show you alerts that are relevant to you.

There are six notification types. Reply notifications fire when someone responds
to a message thread you started. Reaction notifications fire when someone reacts
with an emoji to your message. Pin notifications fire when a channel moderator
pins an important message. Task notifications fire when a Kanban card is
assigned to you. Mention notifications fire instantly when someone types your
@name in a channel. And Meeting notifications fire when an instant or scheduled
meeting is created — and these uniquely include a "Join Meeting" button that
opens the Jitsi room directly.

Each notification type has its own color, icon, and gradient style so you can
identify it at a glance without reading the text.

The notification panel filters by your email address. Even if a broadcast
notification goes to the whole org, you only see the ones where your email is
in the recipient list. This prevents notification overload.

The presence system updates in real time. When you log in, your dot turns green.
After five minutes of no interaction, it shifts to amber — away. If you manually
set Do Not Disturb in your settings, it turns red and suppresses non-critical
notifications. When you close the browser or log out, it turns grey within 30
seconds via the heartbeat timeout.

============================================================
SLIDE 15 - PROS, CONS & HONEST ASSESSMENT (3 minutes)
============================================================

I want to be transparent about EduTechExOS — both what it does brilliantly
and where there are genuine limitations.

The pros are significant.

The biggest one is cost. Zero per-seat SaaS cost. We pay for compute, not for
seats. For a team of ten, we're saving tens of thousands of rupees per year
compared to Slack plus Notion plus Zoom.

Complete customizability. Want a new feature? We build it. No feature request
queue, no waiting for a vendor's quarterly release.

Deep AI integration. Not a chatbot widget — AI that knows your workspace.

End-to-end AES-256 message encryption. Your conversations are private even from
us as platform operators.

Jitsi meetings are free and have no participant limits. We could have 100 people
in a call at no cost.

Single application replacing five separate tools.

Now for the honest cons.

The Render free tier has a 50-second cold start. The first person to visit each
morning triggers a cold boot. We handle this with a retry mechanism, but it's
still a rough experience. A paid Render tier would eliminate this.

No native mobile app yet. You access it through a mobile browser, which works
but isn't as polished as a dedicated app.

File size limits: five megabytes for raw files, ten for media. For teams that
share large video recordings or design files, this needs improvement.

B2 and R2 storage aren't configured yet, so the fallback is MongoDB base64 for
raw files, which has storage implications at scale.

No offline mode. No calendar sync with external calendars. Single tenant.
Limited third-party integrations.

These are all solvable — they're roadmap items, not fundamental limitations.

============================================================
SLIDE 16 - WHY ONLY THIS APPLICATION (4 minutes)
============================================================

This is the question I want to answer definitively: why not just use Slack,
or Teams, or Notion?

The first reason is that EduTechExOS was built FOR us BY us. Every single
feature on this platform exists because someone on the EduTechEx team needed it.
There are no features we never use. There are no missing features we have to
work around. When we need something new, we build it in hours, not months.

The second reason is cost. Let me be specific. Slack Pro costs around 7.25
dollars per seat per month. For 10 people that's 72 dollars per month — about
six thousand rupees. Add Notion at 8 dollars per user and Zoom at 14 dollars
per host, and you're easily at 15,000 rupees per month for tools alone. Over
a year that's 180,000 rupees. EduTechExOS runs on Render and Vercel free tiers
— effective cost is near zero.

The third reason is data ownership. When you use Slack, Slack reads your
messages for AI features, compliance, and product improvement. With EduTechExOS,
your data lives in your MongoDB instance. Nobody else reads it. With AES-256
encryption at rest, even a database breach exposes nothing readable.

The fourth reason is AI context. Generic AI tools like ChatGPT know nothing
about your team. Our Gemini Copilot has read every message in every channel
you've participated in. It knows who your teammates are, what tasks are open,
and what was discussed yesterday. The answers are relevant because the context
is real.

The fifth reason is the single-login experience. Your team members don't have
to context-switch between five applications. Everything — communication, tasks,
knowledge, meetings, HR — is in one tab. Cognitive overhead drops dramatically.

And the sixth reason is growth potential. Because we own the codebase, we can
add a mobile app, calendar sync, LMS integration, or anything else we need.
We are not dependent on any vendor's roadmap.

============================================================
SLIDE 17 - FUTURE ROADMAP & THANK YOU (3 minutes)
============================================================

Let me close by showing you where EduTechExOS is going.

In Q3 2026, we plan to ship a native mobile app using React Native. iOS and
Android, with push notifications, so the team is connected even away from
their desks.

Also in Q3, we'll add Google Calendar and Outlook sync. Scheduled meetings
and approved leave requests will automatically appear in your personal calendar.

In Q4 2026, we'll fully configure Backblaze B2 storage, lifting the file size
restrictions for raw uploads. We'll also add multi-workspace support — one
platform instance serving multiple independent teams within EduTechEx.

In Q1 2027, we'll integrate with the LMS — linking course materials and
assessments directly to channel discussions. We'll also add AI-generated
meeting minutes: after a recorded meeting ends, Gemini transcribes and
summarizes it automatically.

In Q2 2027, a public API and Zapier integration will let external tools trigger
EduTechExOS actions — for example, automatically creating a Kanban card when a
support ticket is filed. And we'll implement end-to-end encrypted direct
messages using the Signal protocol.

This is a living platform. It evolves with us.

I want to close by saying: EduTechExOS is not just software. It's infrastructure
for how this team thinks, communicates, and works together. Every feature we've
built — the AI that knows your workspace, the attendance that tracks itself,
the meetings that anyone can join with one click — was built with one goal:
to make your workday more focused, more connected, and more productive.

Thank you very much. I'm happy to take questions.

============================================================
TOTAL ESTIMATED SPEAKING TIME: 45-55 minutes
============================================================

Tips for Q&A:
- If asked about cost: emphasize zero per-seat SaaS licensing
- If asked about security: AES-256, bcrypt, JWT, single-use invite tokens
- If asked about mobile: React Native app is Q3 2026 roadmap
- If asked about file limits: mention B2 integration on Q4 2026 roadmap
- If asked about offline: currently not supported; PWA is a possible future direction
"""

speech_path = "EduTechExOS_Speech.txt"
with open(speech_path, "w", encoding="utf-8") as f:
    f.write(speech.strip())
print(f"[OK] Saved {speech_path}")
print("Done.")
